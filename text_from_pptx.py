#!/usr/bin/env python3
import sys

from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    """
    Extracts all text from a PowerPoint presentation (.pptx).
    """
    try:
        presentation = Presentation(pptx_file)
        full_text = []
        for slide_number, slide in enumerate(presentation.slides):
            full_text.append(f"\n--- Slide {slide_number + 1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"An error occurred: {e}"


def main():
    if len(sys.argv) < 2:
        print('You must supply a pptx to extract from.')
        exit(1)
    file_path = sys.argv[1]
    extracted_text = extract_text_from_pptx(file_path)
    print(extracted_text)


if __name__ == '__main__':
    main()

